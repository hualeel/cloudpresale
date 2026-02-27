"""
PPT 生成器 (python-pptx)
- generate_ppt_overview: 技术方案 PPT（详细版，面向技术团队）
- generate_ppt_exec: 汇报方案 PPT（精简版，面向决策层）
"""
import io
import re
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── 设计常量 ─────────────────────────────────────────
BG_DARK    = RGBColor(0x08, 0x0D, 0x18)   # 深色背景
COLOR_ACC1 = RGBColor(0x0A, 0x84, 0xFF)   # 主色蓝
COLOR_ACC2 = RGBColor(0x30, 0xD1, 0x58)   # 强调绿
COLOR_TEXT = RGBColor(0xE2, 0xE8, 0xF0)   # 主文字
COLOR_GRAY = RGBColor(0x94, 0xA3, 0xB8)   # 辅助灰

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    """使用空白布局创建新幻灯片。"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _set_bg(slide, color: RGBColor = BG_DARK):
    """设置幻灯片背景颜色。"""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text: str, left, top, width, height,
                   font_size=18, bold=False, color=COLOR_TEXT,
                   align=PP_ALIGN.LEFT, wrap=True):
    """添加文本框，返回文本框对象。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _add_title_slide(prs: Presentation, title: str, subtitle: str, customer: str, date_str: str):
    """封面幻灯片。"""
    slide = _blank_slide(prs)
    _set_bg(slide)

    # 装饰线（顶部）
    from pptx.util import Emu
    from pptx.dml.color import RGBColor
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_W, Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACC1
    line.line.fill.background()

    # 主标题
    _add_text_box(slide, title,
                  Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                  font_size=36, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER)

    # 副标题（客户名）
    _add_text_box(slide, f"— {customer} —",
                  Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
                  font_size=20, bold=False, color=COLOR_ACC1, align=PP_ALIGN.CENTER)

    # 日期
    _add_text_box(slide, date_str,
                  Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
                  font_size=14, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

    # 底部装饰线
    line2 = slide.shapes.add_shape(
        1, Inches(0), Inches(7.45), SLIDE_W, Inches(0.05)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = COLOR_ACC2
    line2.line.fill.background()


def _add_section_slide(prs: Presentation, section_title: str, section_num: str):
    """章节过渡页。"""
    slide = _blank_slide(prs)
    _set_bg(slide)

    # 编号
    _add_text_box(slide, section_num,
                  Inches(1), Inches(2.5), Inches(1.5), Inches(1.5),
                  font_size=64, bold=True, color=COLOR_ACC1, align=PP_ALIGN.LEFT)

    # 标题
    _add_text_box(slide, section_title,
                  Inches(2.5), Inches(3.0), Inches(9), Inches(1.2),
                  font_size=32, bold=True, color=COLOR_TEXT, align=PP_ALIGN.LEFT)

    # 分隔线
    rect = slide.shapes.add_shape(1, Inches(2.5), Inches(4.2), Inches(6), Inches(0.04))
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_ACC1
    rect.line.fill.background()


def _add_content_slide(prs: Presentation, title: str, content_lines: list[str]):
    """内容页：左侧标题 + 右侧正文。"""
    slide = _blank_slide(prs)
    _set_bg(slide)

    # 标题条
    rect = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x0D, 0x1A, 0x35)
    rect.line.fill.background()

    _add_text_box(slide, title,
                  Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
                  font_size=20, bold=True, color=COLOR_TEXT)

    # 正文区域
    content_text = "\n".join(content_lines[:18])  # 最多18行
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines[:18]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if line.startswith('**') and line.endswith('**'):
            # 粗体行
            run = p.add_run()
            run.text = line.strip('*')
            run.font.bold = True
            run.font.color.rgb = COLOR_ACC1
            run.font.size = Pt(13)
        elif line.startswith('• ') or line.startswith('- '):
            run = p.add_run()
            run.text = "  " + line
            run.font.color.rgb = COLOR_TEXT
            run.font.size = Pt(12)
            p.space_before = Pt(4)
        else:
            run = p.add_run()
            run.text = line
            run.font.color.rgb = COLOR_GRAY
            run.font.size = Pt(11)


def _markdown_to_slides(prs: Presentation, section_title: str, section_num: str, content: str):
    """
    将一个 Agent 的 Markdown 内容转换为 1个章节页 + 多个内容页。
    """
    _add_section_slide(prs, section_title, section_num)

    if not content:
        return

    # 按 ## 子标题分割为多页
    sub_sections = re.split(r'\n(?=## )', content.strip())

    for sub in sub_sections:
        lines = sub.strip().split('\n')
        if not lines:
            continue

        # 提取子标题
        if lines[0].startswith('## '):
            slide_title = lines[0][3:]
            body_lines = lines[1:]
        elif lines[0].startswith('# '):
            slide_title = lines[0][2:]
            body_lines = lines[1:]
        else:
            slide_title = section_title
            body_lines = lines

        # 清洗正文行
        clean_lines = []
        for line in body_lines:
            line = line.rstrip()
            if line.startswith('### '):
                clean_lines.append(f"**{line[4:]}**")
            elif line.startswith('#### '):
                clean_lines.append(f"**{line[5:]}**")
            elif re.match(r'^\|[-:\s|]+\|$', line):
                continue  # 跳过表格分隔行
            elif line.startswith('|'):
                # 表格行 → 转为列表
                cells = [c.strip() for c in line.strip('|').split('|')]
                clean_lines.append("• " + " | ".join(cells))
            elif line.startswith('- ') or line.startswith('* '):
                clean_lines.append("• " + line[2:])
            elif line == '':
                if clean_lines and clean_lines[-1] != '':
                    clean_lines.append('')
            else:
                clean_lines.append(line)

        # 每页最多 15 行，超出分页
        page_size = 15
        for page_start in range(0, max(1, len(clean_lines)), page_size):
            page_lines = clean_lines[page_start:page_start + page_size]
            page_title = slide_title if page_start == 0 else f"{slide_title}（续）"
            _add_content_slide(prs, page_title, page_lines)


# ── 公开接口 ────────────────────────────────────────

def generate_ppt_overview(
    solution_content: dict,
    requirement: dict,
    customer_name: str,
    version: str,
) -> bytes:
    """
    生成技术方案 PPT（完整版，面向技术团队）。
    包含全部 6 个 Agent 章节。
    """
    prs = _new_prs()
    date_str = datetime.now().strftime("%Y年%m月%d日")
    doc_title = requirement.get("title", "云原生技术方案")

    _add_title_slide(prs, f"{customer_name}\n{doc_title}", "技术方案详述", customer_name, date_str)

    SECTIONS = [
        ("arch",      "架构设计方案", "01"),
        ("sizing",    "资源规格设计", "02"),
        ("security",  "安全合规方案", "03"),
        ("migration", "应用迁移路径", "04"),
        ("plan",      "项目实施计划", "05"),
        ("pricing",   "商务报价方案", "06"),
    ]

    for key, section_title, num in SECTIONS:
        agent_result = solution_content.get(key)
        if not agent_result:
            continue
        content_text = agent_result.get("content", "") if isinstance(agent_result, dict) else str(agent_result)
        _markdown_to_slides(prs, section_title, num, content_text)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def generate_ppt_exec(
    solution_content: dict,
    requirement: dict,
    customer_name: str,
    version: str,
) -> bytes:
    """
    生成汇报方案 PPT（精简版，面向决策层）。
    只包含：架构总览、资源摘要、安全合规、商务报价，共约 8-12 页。
    """
    prs = _new_prs()
    date_str = datetime.now().strftime("%Y年%m月%d日")
    doc_title = requirement.get("title", "云原生建设方案")

    _add_title_slide(prs, f"{customer_name}\n{doc_title}", "汇报方案", customer_name, date_str)

    # 执行摘要页
    exec_lines = _build_exec_summary(solution_content, requirement, customer_name)
    _add_content_slide(prs, "执行摘要", exec_lines)

    # 精简版章节（只取每个Agent前两个子章节）
    EXEC_SECTIONS = [
        ("arch",     "方案架构总览", "01"),
        ("security", "安全合规要点", "02"),
        ("plan",     "实施路线图",  "03"),
        ("pricing",  "投资与回报",  "04"),
    ]

    for key, section_title, num in EXEC_SECTIONS:
        agent_result = solution_content.get(key)
        if not agent_result:
            continue
        content_text = agent_result.get("content", "") if isinstance(agent_result, dict) else str(agent_result)
        # 只取前600字符（精简版）
        short_content = content_text[:800] if len(content_text) > 800 else content_text
        _markdown_to_slides(prs, section_title, num, short_content)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _build_exec_summary(solution_content: dict, requirement: dict, customer_name: str) -> list[str]:
    """构建执行摘要页的要点列表。"""
    lines = [
        f"**客户：{customer_name}**",
        "",
        f"• 项目背景：{requirement.get('content', {}).get('pain_points', '业务数字化转型需求')}",
        f"• 目标容器化比例：{requirement.get('content', {}).get('target_containerization', '≥80%')}",
        f"• 合规要求：{'、'.join(requirement.get('content', {}).get('compliance', ['等保三级']))}",
        f"• 预算范围：{requirement.get('content', {}).get('budget_range', '待确认')}",
        "",
        "**核心价值主张：**",
        "• 一站式云原生平台，覆盖容器、DevOps、微服务全栈",
        "• 金融行业最佳实践，满足监管合规要求",
        "• 专业实施服务，保障平滑迁移",
    ]
    return lines


# ── 专项 PPT ─────────────────────────────────────────

def generate_ppt_container(
    solution_content: dict,
    requirement: dict,
    customer_name: str,
    version: str,
) -> bytes:
    """
    生成容器平台专项 PPT（面向容器平台技术团队 / IT 基础架构部门）。
    聚焦：架构设计（全量）+ 资源规格（全量）+ 迁移路径（摘要）。
    """
    prs = _new_prs()
    date_str = datetime.now().strftime("%Y年%m月%d日")
    doc_title = requirement.get("title", "容器平台建设方案")

    _add_title_slide(prs, f"{customer_name}\n{doc_title}", "容器平台专项", customer_name, date_str)

    # 平台建设背景与目标摘要页
    _add_content_slide(prs, "平台建设背景与目标",
                       _build_container_summary(solution_content, requirement, customer_name))

    arch = solution_content.get("arch")
    if arch:
        text = arch.get("content", "") if isinstance(arch, dict) else str(arch)
        _markdown_to_slides(prs, "容器平台架构设计", "01", text)

    sizing = solution_content.get("sizing")
    if sizing:
        text = sizing.get("content", "") if isinstance(sizing, dict) else str(sizing)
        _markdown_to_slides(prs, "集群资源规划", "02", text)

    migration = solution_content.get("migration")
    if migration:
        text = migration.get("content", "") if isinstance(migration, dict) else str(migration)
        _markdown_to_slides(prs, "应用上容器平台路径", "03", text[:800])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _build_container_summary(solution_content: dict, requirement: dict, customer_name: str) -> list[str]:
    content = requirement.get("content", {})
    return [
        f"**客户：{customer_name}**",
        "",
        f"• 当前容器化比例：{content.get('current_containerization', '待确认')}",
        f"• 目标容器化比例：{content.get('target_containerization', '≥80%')}",
        f"• 集群数量规划：{content.get('cluster_count', '待确认')}",
        f"• 核心建设模块：{'、'.join(content.get('modules', ['容器平台', 'DevOps', '微服务']))}",
        "",
        "**建设目标：**",
        "• 建立统一云原生底座，支撑全行应用容器化改造",
        "• 实现 CI/CD 流水线自动化，提升研发效能",
        "• 满足金融监管合规要求，保障平台安全稳定运行",
    ]


def generate_ppt_devops(
    solution_content: dict,
    requirement: dict,
    customer_name: str,
    version: str,
) -> bytes:
    """
    生成 DevOps 专项 PPT（面向开发团队 / DevOps 工程师）。
    聚焦：应用迁移策略（全量）+ 实施计划（全量）+ 架构 CI/CD 部分（摘要）。
    """
    prs = _new_prs()
    date_str = datetime.now().strftime("%Y年%m月%d日")
    doc_title = requirement.get("title", "DevOps 转型方案")

    _add_title_slide(prs, f"{customer_name}\n{doc_title}", "DevOps 专项", customer_name, date_str)

    _add_content_slide(prs, "DevOps 转型目标",
                       _build_devops_summary(solution_content, requirement, customer_name))

    migration = solution_content.get("migration")
    if migration:
        text = migration.get("content", "") if isinstance(migration, dict) else str(migration)
        _markdown_to_slides(prs, "应用迁移与上云路径", "01", text)

    plan = solution_content.get("plan")
    if plan:
        text = plan.get("content", "") if isinstance(plan, dict) else str(plan)
        _markdown_to_slides(prs, "项目实施计划", "02", text)

    arch = solution_content.get("arch")
    if arch:
        text = arch.get("content", "") if isinstance(arch, dict) else str(arch)
        _markdown_to_slides(prs, "CI/CD 与平台架构", "03", text[:800])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _build_devops_summary(solution_content: dict, requirement: dict, customer_name: str) -> list[str]:
    content = requirement.get("content", {})
    return [
        f"**客户：{customer_name}**",
        "",
        f"• 目标容器化比例：{content.get('target_containerization', '≥80%')}",
        f"• 预算范围：{content.get('budget_range', '待确认')}",
        f"• 核心痛点：{content.get('pain_points', '研发效能低、环境一致性差、发布风险高')}",
        "",
        "**DevOps 转型目标：**",
        "• 建立标准化 CI/CD 流水线，实现一键发布",
        "• 统一容器镜像管理，保障镜像供应链安全",
        "• 实现蓝绿发布 / 金丝雀发布，降低发布风险",
        "• 建立研发效能度量体系（DORA Metrics）",
    ]


def generate_ppt_security(
    solution_content: dict,
    requirement: dict,
    customer_name: str,
    version: str,
) -> bytes:
    """
    生成安全合规专项 PPT（面向安全团队 / 合规部门 / CISO）。
    聚焦：安全合规方案（全量）+ 架构安全要点（摘要）。
    """
    prs = _new_prs()
    date_str = datetime.now().strftime("%Y年%m月%d日")
    doc_title = requirement.get("title", "云原生安全合规方案")

    _add_title_slide(prs, f"{customer_name}\n{doc_title}", "安全合规专项", customer_name, date_str)

    _add_content_slide(prs, "合规框架与安全目标",
                       _build_security_summary(solution_content, requirement, customer_name))

    security = solution_content.get("security")
    if security:
        text = security.get("content", "") if isinstance(security, dict) else str(security)
        _markdown_to_slides(prs, "安全合规方案详述", "01", text)

    arch = solution_content.get("arch")
    if arch:
        text = arch.get("content", "") if isinstance(arch, dict) else str(arch)
        _markdown_to_slides(prs, "安全架构设计要点", "02", text[:800])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _build_security_summary(solution_content: dict, requirement: dict, customer_name: str) -> list[str]:
    content = requirement.get("content", {})
    compliance_list = content.get("compliance", ["等保三级"])
    return [
        f"**客户：{customer_name}**",
        "",
        f"• 合规要求：{'、'.join(compliance_list)}",
        "• 行业监管：银保监会 / 人民银行 / 证监会",
        "• 安全目标：满足金融行业网络安全等级保护标准",
        "",
        "**安全建设范围：**",
        "• 身份认证与访问控制（RBAC / MFA）",
        "• 网络微隔离与 API 网关安全",
        "• 镜像扫描与供应链安全",
        "• 运行时威胁检测与审计日志",
        "• 数据传输与存储加密",
    ]
